from unittest.mock import patch

from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

from app.services.document_parser import (
    media_type_for_filename,
    parse_document,
    safe_filename,
)
from app.utils.chunker import chunk_document_blocks


def test_filename_allowlist_and_sanitization():
    assert safe_filename("../../draft?.docx") == "draft_.docx"
    assert media_type_for_filename("paper.PDF") == "application/pdf"


def test_parse_html_preserves_text_and_table_modalities(tmp_path):
    path = tmp_path / "report.html"
    path.write_text(
        """
        <html><head><title>Experiment</title><script>ignore()</script></head>
        <body><h1>Results</h1><p>The model improved retrieval quality.</p>
        <table><tr><th>Metric</th><th>Value</th></tr>
        <tr><td>Recall@5</td><td>0.75</td></tr></table></body></html>
        """,
        encoding="utf-8",
    )

    parsed = parse_document(path)

    assert parsed.title == "Experiment"
    assert {block.modality for block in parsed.blocks} == {"text", "table"}
    assert "ignore()" not in "\n".join(block.text for block in parsed.blocks)
    assert parsed.metadata["modalities"] == ["table", "text"]


def test_parse_docx_and_preserve_table_locator(tmp_path):
    path = tmp_path / "paper.docx"
    document = DocxDocument()
    document.core_properties.title = "Paper Notes"
    document.add_heading("Method", level=1)
    document.add_paragraph(
        "Hybrid retrieval combines dense candidates with BM25 lexical evidence."
    )
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Recall@5"
    table.cell(1, 1).text = "75.31"
    document.save(path)

    parsed = parse_document(path)

    assert parsed.title == "Paper Notes"
    table_block = next(block for block in parsed.blocks if block.modality == "table")
    assert table_block.source_locator == {"table": 1}
    assert "Recall@5" in table_block.text


def test_parse_pptx_and_xlsx_with_structural_locators(tmp_path):
    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Architecture"
    slide.placeholders[1].text = "Eight control nodes and thirteen visible stages"
    presentation.save(pptx_path)

    parsed_pptx = parse_document(pptx_path)
    assert parsed_pptx.metadata["slide_count"] == 1
    assert parsed_pptx.blocks[0].source_locator["slide"] == 1

    xlsx_path = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Eval"
    sheet.append(["Metric", "Dense", "Hybrid"])
    sheet.append(["Recall@5", 73.70, 75.31])
    workbook.save(xlsx_path)

    parsed_xlsx = parse_document(xlsx_path)
    assert parsed_xlsx.blocks[0].modality == "table"
    assert parsed_xlsx.blocks[0].source_locator == {"sheet": "Eval"}


def test_image_ocr_and_chunk_provenance(tmp_path):
    path = tmp_path / "figure.png"
    Image.new("RGB", (640, 240), "white").save(path)

    with patch(
        "app.services.document_parser.pytesseract.image_to_string",
        return_value="Evidence pipeline diagram with citation gate.",
    ):
        parsed = parse_document(path)

    chunks = chunk_document_blocks(parsed.blocks)
    assert parsed.blocks[0].modality == "image_ocr"
    assert chunks[0].modality == "image_ocr"
    assert chunks[0].source_locator == {"image": 1}
