"""Normalize heterogeneous files into provenance-preserving document blocks.

Supported inputs:

- PDF, including text-poor/scanned pages and embedded-image OCR
- DOCX and PPTX, including tables and embedded-image OCR
- HTML, Markdown, plain text, and CSV
- XLSX workbooks (one table block per sheet)
- PNG/JPEG/WebP/TIFF images through local Tesseract OCR

Parsing is deliberately local and deterministic. OCR failures are recorded as
warnings; a document only succeeds when at least one usable text/table/OCR
block is produced.
"""
from __future__ import annotations

import csv
import io
import re
import zipfile
from pathlib import Path
from typing import Iterable

import pdfplumber
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    import fitz  # type: ignore
except Exception:  # pragma: no cover - requirements install PyMuPDF
    fitz = None  # type: ignore

try:
    import pytesseract
except Exception:  # pragma: no cover - requirements install pytesseract
    pytesseract = None  # type: ignore

from app.core.config import get_settings
from app.utils.document import DocumentBlock, ParsedDocument


class UnsupportedDocumentError(ValueError):
    """Raised when a file extension is outside the ingestion allow-list."""


class UnsafeArchiveError(ValueError):
    """Raised when an Office archive exceeds configured expansion limits."""


SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".html": "text/html",
    ".htm": "text/html",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".tif": "image/tiff",
    ".tiff": "image/tiff",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}


def safe_filename(raw_name: str | None) -> str:
    """Drop path traversal and normalize problematic filename characters."""
    name = Path(raw_name or "upload").name.strip()
    name = re.sub(r"[\x00-\x1f<>:\"/\\\\|?*]+", "_", name)
    return name[:240] or "upload"


def media_type_for_filename(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    try:
        return SUPPORTED_EXTENSIONS[suffix]
    except KeyError as exc:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise UnsupportedDocumentError(
            f"unsupported file extension {suffix or '(none)'}; supported: {supported}"
        ) from exc


def _warn_once(warnings: list[str], message: str) -> None:
    if message not in warnings:
        warnings.append(message)


def _validate_office_archive(path: Path) -> None:
    settings = get_settings()
    max_bytes = settings.ingest_max_archive_uncompressed_mb * 1024 * 1024
    with zipfile.ZipFile(path) as archive:
        total = sum(max(0, item.file_size) for item in archive.infolist())
        if total > max_bytes:
            raise UnsafeArchiveError(
                f"archive expands to {total} bytes; limit is {max_bytes} bytes"
            )
        for item in archive.infolist():
            if item.file_size > 0 and item.compress_size > 0:
                ratio = item.file_size / item.compress_size
                if ratio > 200:
                    raise UnsafeArchiveError(
                        f"archive member has unsafe compression ratio: {item.filename}"
                    )


def _table_text(rows: Iterable[Iterable[object]]) -> str:
    lines: list[str] = []
    for row in rows:
        values = [str(value).strip() if value is not None else "" for value in row]
        if any(values):
            lines.append(" | ".join(values))
    return "\n".join(lines).strip()


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _ocr_image(
    image: Image.Image,
    *,
    warnings: list[str],
    locator: str,
) -> str:
    settings = get_settings()
    if not settings.ingest_ocr_enabled:
        _warn_once(warnings, "ocr_disabled")
        return ""
    if pytesseract is None:
        _warn_once(warnings, "pytesseract_not_installed")
        return ""
    try:
        return (
            pytesseract.image_to_string(
                image.convert("RGB"),
                lang=settings.ingest_ocr_languages,
            )
            or ""
        ).strip()
    except Exception as exc:
        _warn_once(warnings, f"ocr_failed:{locator}:{type(exc).__name__}")
        return ""


def _parse_pdf(path: Path) -> ParsedDocument:
    warnings: list[str] = []
    blocks: list[DocumentBlock] = []
    settings = get_settings()

    if fitz is None:
        raise RuntimeError("PyMuPDF is required for PDF multimodal parsing")

    document = fitz.open(path)
    page_count = len(document)
    try:
        for page_index, page in enumerate(document, start=1):
            text = (page.get_text("text") or "").strip()
            if text:
                blocks.append(
                    DocumentBlock(
                        text=text,
                        modality="text",
                        source_locator={"page": page_index},
                    )
                )

            # OCR the rendered page when native extraction is sparse. This
            # covers scanned pages without duplicating text-rich pages.
            if len(text) < settings.ingest_ocr_min_text_chars:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                with Image.open(io.BytesIO(pixmap.tobytes("png"))) as image:
                    ocr_text = _ocr_image(
                        image,
                        warnings=warnings,
                        locator=f"page:{page_index}",
                    )
                if ocr_text and ocr_text not in text:
                    blocks.append(
                        DocumentBlock(
                            text=ocr_text,
                            modality="image_ocr",
                            source_locator={"page": page_index, "kind": "page_render"},
                        )
                    )
                continue

            # Text-rich pages can still contain information-bearing figures.
            seen_xrefs: set[int] = set()
            for image_index, image_info in enumerate(page.get_images(full=True), start=1):
                xref = int(image_info[0])
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                extracted = document.extract_image(xref)
                width = int(extracted.get("width") or 0)
                height = int(extracted.get("height") or 0)
                raw = extracted.get("image") or b""
                if width < 160 or height < 100 or len(raw) < 2048:
                    continue
                try:
                    with Image.open(io.BytesIO(raw)) as image:
                        ocr_text = _ocr_image(
                            image,
                            warnings=warnings,
                            locator=f"page:{page_index}:image:{image_index}",
                        )
                except Exception as exc:
                    _warn_once(
                        warnings,
                        f"image_decode_failed:page:{page_index}:{type(exc).__name__}",
                    )
                    continue
                if ocr_text:
                    blocks.append(
                        DocumentBlock(
                            text=ocr_text,
                            modality="image_ocr",
                            source_locator={
                                "page": page_index,
                                "image": image_index,
                                "kind": "embedded_image",
                            },
                        )
                    )
    finally:
        document.close()

    # Preserve table modality separately from the page text so retrieval can
    # explain whether evidence came from prose, a table, or OCR.
    try:
        with pdfplumber.open(path) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                for table_index, table in enumerate(page.extract_tables() or [], start=1):
                    table_text = _table_text(table)
                    if table_text:
                        blocks.append(
                            DocumentBlock(
                                text=table_text,
                                modality="table",
                                source_locator={
                                    "page": page_index,
                                    "table": table_index,
                                },
                            )
                        )
    except Exception as exc:
        _warn_once(warnings, f"pdf_table_extraction_failed:{type(exc).__name__}")

    return ParsedDocument(
        blocks=blocks,
        title=path.stem,
        warnings=warnings,
        metadata={"page_count": page_count},
    )


def _parse_docx(path: Path) -> ParsedDocument:
    _validate_office_archive(path)
    warnings: list[str] = []
    blocks: list[DocumentBlock] = []
    document = DocxDocument(path)
    current_section: str | None = None
    section_parts: list[str] = []
    section_index = 0

    def flush_section() -> None:
        nonlocal section_parts, section_index
        text = "\n".join(section_parts).strip()
        if text:
            section_index += 1
            blocks.append(
                DocumentBlock(
                    text=text,
                    modality="text",
                    source_locator={"section_index": section_index},
                    section=current_section,
                )
            )
        section_parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style_name.startswith("heading"):
            flush_section()
            current_section = text
        else:
            section_parts.append(text)
    flush_section()

    for table_index, table in enumerate(document.tables, start=1):
        text = _table_text([[cell.text for cell in row.cells] for row in table.rows])
        if text:
            blocks.append(
                DocumentBlock(
                    text=text,
                    modality="table",
                    source_locator={"table": table_index},
                    section=current_section,
                )
            )

    image_index = 0
    for relationship in document.part.rels.values():
        target = getattr(relationship, "target_part", None)
        content_type = str(getattr(target, "content_type", ""))
        if not content_type.startswith("image/"):
            continue
        image_index += 1
        try:
            with Image.open(io.BytesIO(target.blob)) as image:
                text = _ocr_image(
                    image,
                    warnings=warnings,
                    locator=f"image:{image_index}",
                )
        except Exception as exc:
            _warn_once(warnings, f"image_decode_failed:{image_index}:{type(exc).__name__}")
            continue
        if text:
            blocks.append(
                DocumentBlock(
                    text=text,
                    modality="image_ocr",
                    source_locator={"image": image_index},
                )
            )

    return ParsedDocument(
        blocks=blocks,
        title=document.core_properties.title or path.stem,
        warnings=warnings,
    )


def _parse_pptx(path: Path) -> ParsedDocument:
    _validate_office_archive(path)
    warnings: list[str] = []
    blocks: list[DocumentBlock] = []
    presentation = Presentation(path)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    slide_text.append(text)
            if getattr(shape, "has_table", False):
                text = _table_text(
                    [[cell.text for cell in row.cells] for row in shape.table.rows]
                )
                if text:
                    blocks.append(
                        DocumentBlock(
                            text=text,
                            modality="table",
                            source_locator={"slide": slide_index, "shape": shape_index},
                        )
                    )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    with Image.open(io.BytesIO(shape.image.blob)) as image:
                        text = _ocr_image(
                            image,
                            warnings=warnings,
                            locator=f"slide:{slide_index}:shape:{shape_index}",
                        )
                except Exception as exc:
                    _warn_once(
                        warnings,
                        f"image_decode_failed:slide:{slide_index}:{type(exc).__name__}",
                    )
                    continue
                if text:
                    blocks.append(
                        DocumentBlock(
                            text=text,
                            modality="image_ocr",
                            source_locator={"slide": slide_index, "shape": shape_index},
                        )
                    )
        if slide_text:
            blocks.append(
                DocumentBlock(
                    text="\n".join(slide_text),
                    modality="text",
                    source_locator={"slide": slide_index},
                )
            )

    return ParsedDocument(
        blocks=blocks,
        title=presentation.core_properties.title or path.stem,
        warnings=warnings,
        metadata={"slide_count": len(presentation.slides)},
    )


def _parse_html(path: Path) -> ParsedDocument:
    soup = BeautifulSoup(_decode_text(path.read_bytes()), "html.parser")
    for element in soup(["script", "style", "noscript", "template"]):
        element.decompose()

    blocks: list[DocumentBlock] = []
    for table_index, table in enumerate(list(soup.find_all("table")), start=1):
        rows = [
            [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
            for row in table.find_all("tr")
        ]
        text = _table_text(rows)
        if text:
            blocks.append(
                DocumentBlock(
                    text=text,
                    modality="table",
                    source_locator={"table": table_index},
                )
            )
        table.decompose()

    body = soup.body or soup
    text = body.get_text("\n", strip=True)
    if text:
        blocks.insert(
            0,
            DocumentBlock(text=text, modality="text", source_locator={"kind": "html_body"}),
        )
    title = soup.title.get_text(" ", strip=True) if soup.title else path.stem
    return ParsedDocument(blocks=blocks, title=title)


def _parse_text(path: Path, modality: str = "text") -> ParsedDocument:
    text = _decode_text(path.read_bytes()).strip()
    blocks = [
        DocumentBlock(text=text, modality=modality, source_locator={"kind": path.suffix.lower()[1:]})
    ] if text else []
    return ParsedDocument(blocks=blocks, title=path.stem)


def _parse_csv(path: Path) -> ParsedDocument:
    raw = _decode_text(path.read_bytes())
    sample = raw[:8192]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows = list(csv.reader(io.StringIO(raw), dialect))
    text = _table_text(rows)
    blocks = [
        DocumentBlock(text=text, modality="table", source_locator={"sheet": path.stem})
    ] if text else []
    return ParsedDocument(blocks=blocks, title=path.stem, metadata={"row_count": len(rows)})


def _parse_xlsx(path: Path) -> ParsedDocument:
    _validate_office_archive(path)
    blocks: list[DocumentBlock] = []
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            text = _table_text(sheet.iter_rows(values_only=True))
            if text:
                blocks.append(
                    DocumentBlock(
                        text=text,
                        modality="table",
                        source_locator={"sheet": sheet.title},
                        section=sheet.title,
                    )
                )
    finally:
        workbook.close()
    return ParsedDocument(
        blocks=blocks,
        title=path.stem,
        metadata={"sheet_count": len(blocks)},
    )


def _parse_image(path: Path) -> ParsedDocument:
    warnings: list[str] = []
    with Image.open(path) as image:
        width, height = image.size
        text = _ocr_image(image, warnings=warnings, locator="image:1")
    blocks = [
        DocumentBlock(
            text=text,
            modality="image_ocr",
            source_locator={"image": 1},
        )
    ] if text else []
    return ParsedDocument(
        blocks=blocks,
        title=path.stem,
        warnings=warnings,
        metadata={"width": width, "height": height},
    )


def parse_document(path: str | Path) -> ParsedDocument:
    """Parse one allow-listed file and require at least one usable block."""
    resolved = Path(path)
    if not resolved.exists() or not resolved.is_file():
        raise FileNotFoundError(resolved)
    suffix = resolved.suffix.lower()
    media_type_for_filename(resolved.name)

    if suffix == ".pdf":
        parsed = _parse_pdf(resolved)
    elif suffix == ".docx":
        parsed = _parse_docx(resolved)
    elif suffix == ".pptx":
        parsed = _parse_pptx(resolved)
    elif suffix == ".xlsx":
        parsed = _parse_xlsx(resolved)
    elif suffix in {".html", ".htm"}:
        parsed = _parse_html(resolved)
    elif suffix in {".md", ".markdown", ".txt"}:
        parsed = _parse_text(resolved)
    elif suffix == ".csv":
        parsed = _parse_csv(resolved)
    elif suffix in IMAGE_EXTENSIONS:
        parsed = _parse_image(resolved)
    else:  # protected by media_type_for_filename; keeps type check exhaustive
        raise UnsupportedDocumentError(suffix)

    usable = [block for block in parsed.blocks if block.text.strip()]
    if not usable:
        detail = ", ".join(parsed.warnings) or "no text, table, or OCR content found"
        raise RuntimeError(f"no usable document blocks: {detail}")
    return ParsedDocument(
        blocks=usable,
        title=parsed.title,
        warnings=parsed.warnings,
        metadata={
            **parsed.metadata,
            "modalities": sorted({block.modality for block in usable}),
            "block_count": len(usable),
        },
    )


__all__ = [
    "IMAGE_EXTENSIONS",
    "OFFICE_EXTENSIONS",
    "SUPPORTED_EXTENSIONS",
    "UnsafeArchiveError",
    "UnsupportedDocumentError",
    "media_type_for_filename",
    "parse_document",
    "safe_filename",
]
